from pptx import Presentation
from pptx.util import Inches

# Create a Presentation object
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]  # Use the title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Welcome to PPTX Creation Tool"
subtitle.text = "This is a simple example"

# Add a content slide
slide_layout = prs.slide_layouts[1]  # Use the title and content slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "First Page"
content.text = "This is the content of the first page"

# Save the PPTX file
prs.save("example.pptx")
